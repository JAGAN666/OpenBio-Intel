"""FastAPI surface for the LangGraph research agent.

    uvicorn api:app --reload --port 8000

Exposes POST /api/research -> SmartTableResponse, the same validated Pydantic
object the CLI produces. The response_model is that Pydantic class, so FastAPI
also publishes it as the OpenAPI schema at /docs -- backend and frontend share
one contract definition.

Also exposes POST /api/research/stream -> text/event-stream (SSE). The
Map-Reduce extraction stage can fan out to dozens (sometimes hundreds, once
the corpus is large and the question names a popular drug the knowledge
graph resolves broadly) of parallel extract_trial workers -- a synchronous
POST /api/research can genuinely run past a 60s ALB idle timeout on exactly
that shape of question, returning a bare 502 with no indication anything
was even working. /stream reports real progress AS THE GRAPH EXECUTES
instead of going silent until everything finishes, and — just as
importantly for the timeout problem — the ALB sees a continuous trickle of
bytes on the connection the whole time, not one long silence, which is what
an idle-timeout actually watches for.

Also exposes POST /api/landscape -> LandscapeMatrix, a competitive
landscape matrix (mechanism/target rows x development-phase columns) for a
therapeutic area, via a separate, simpler retrieve-then-synthesize graph --
see research_agent.py's LandscapeMatrix/make_landscape_graph docstrings.

Also exposes POST /api/catalysts -> CatalystTimeline, a chronological
tracker of upcoming market-moving events (trial data readouts, PDUFA
dates) for a query -- see research_agent.py's CatalystTimeline/
make_catalyst_graph docstrings, including why this pipeline does a live
ClinicalTrials.gov date lookup that neither /api/research nor
/api/landscape need.

AUTH NOTE: auth.py's Depends(get_current_user) gate was removed from every
endpoint below (was on all 6) -- it validates a JWT but this project never
built anything to ISSUE one (no login endpoint, no frontend token storage),
so every real request from the actual frontend was hitting a 401 with no
way to ever succeed. auth.py itself is untouched and still importable for
whenever a real login flow exists to pair it with; this file just doesn't
call it right now. No tenant_id is threaded through the graphs for the same
reason (see _initial_state).
"""

from __future__ import annotations

import io
import json
import logging
import os
import time
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Any, AsyncIterator

from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, StreamingResponse
from langchain_core.messages import HumanMessage
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

from research_agent import (
    COLLECTION_NAME,
    DEFAULT_MODEL,
    CatalystTimeline,
    LandscapeMatrix,
    QDRANT_HOST,
    QDRANT_PORT,
    SmartTableResponse,
    make_catalyst_graph,
    make_graph,
    make_landscape_graph,
    run_catalyst_query,
    run_landscape_query,
)

load_dotenv(Path(__file__).resolve().parent / ".env", override=False)

logging.basicConfig(level=logging.INFO, format="%(asctime)s  %(levelname)-7s %(message)s")
log = logging.getLogger("medical-rag.api")

# The Next.js dev server. Add deployed origins here rather than using "*" --
# a wildcard would let any page on the internet drive this agent.
ALLOWED_ORIGINS = ["http://localhost:3000", "http://127.0.0.1:3000"]

# Compiled once at startup, reused per request. Building the graph constructs a
# ChatAnthropic client and binds tool schemas; doing that per request would add
# avoidable latency to every call. The graph itself is stateless -- all state
# lives in the dict passed to invoke() -- so sharing it is safe.
_graph = None
# Same reasoning, separate graph: the landscape matrix pipeline is a
# distinct, simpler retrieve-then-synthesize shape (see
# make_landscape_graph's own docstring for why it isn't the six-tool ReACT
# agent above), so it gets its own compiled instance rather than being
# forced through _graph's tool-calling loop.
_landscape_graph = None
# Same reasoning again: the catalyst tracker is its own retrieve-then-
# synthesize graph (with an extra live-date-enrichment step neither of the
# other two graphs needs -- see make_catalyst_graph's docstring), so it
# also gets its own compiled instance.
_catalyst_graph = None


@asynccontextmanager
async def lifespan(app: FastAPI):
    global _graph, _landscape_graph, _catalyst_graph
    if not os.getenv("ANTHROPIC_API_KEY"):
        log.warning("ANTHROPIC_API_KEY is not set — /api/research will fail")
    log.info("compiling LangGraph agent (model=%s)…", DEFAULT_MODEL)
    _graph = make_graph(DEFAULT_MODEL, verbose=False)
    log.info("compiling landscape matrix graph (model=%s)…", DEFAULT_MODEL)
    _landscape_graph = make_landscape_graph(DEFAULT_MODEL, verbose=False)
    log.info("compiling catalyst tracker graph…")
    _catalyst_graph = make_catalyst_graph(verbose=False)
    log.info("agent ready; collection=%s", COLLECTION_NAME)
    yield
    log.info("shutting down")


app = FastAPI(
    title="medical-rag research agent",
    description="LangGraph ReACT agent over a Qdrant clinical-trials corpus.",
    version="1.0.0",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)


class ResearchRequest(BaseModel):
    query: str = Field(
        min_length=3,
        max_length=2000,
        description="Natural-language analyst question about the trials corpus.",
        examples=["Compare the mechanisms and sponsors of Phase 3 oncology trials"],
    )


@app.get("/api/health")
def health() -> dict:
    """Cheap readiness probe — does not call the LLM."""
    from qdrant_client import QdrantClient

    try:
        c = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)
        points = c.get_collection(COLLECTION_NAME).points_count
        qdrant_ok = True
    except Exception as exc:  # noqa: BLE001 - surfaced verbatim to the caller
        points, qdrant_ok = None, False
        log.warning("qdrant health check failed: %s", exc)

    return {
        "status": "ok" if (qdrant_ok and _graph is not None) else "degraded",
        "agent_ready": _graph is not None,
        "qdrant_ok": qdrant_ok,
        "collection": COLLECTION_NAME,
        "points": points,
        "model": DEFAULT_MODEL,
        "anthropic_key_set": bool(os.getenv("ANTHROPIC_API_KEY")),
    }


def _initial_state(query: str, tenant_id: str | None = None) -> dict:
    return {
        "messages": [HumanMessage(content=query)],
        "tenant_id": tenant_id,
        "tool_rounds": 0, "result": None,
        "is_in_domain": None, "has_results": None,
        "synthesis_retries": 0, "synthesis_error": None,
        "retrieved_trials": [], "extracted_rows": [], "retrieved_literature": [],
        "retrieved_fda": [],
    }


@app.post("/api/research", response_model=SmartTableResponse)
def research(req: ResearchRequest) -> SmartTableResponse:
    """Run the agent and return the validated Smart Table.

    Declared `def`, not `async def`, on purpose: graph.invoke() is blocking
    (network I/O to Anthropic plus local ONNX embedding), so FastAPI runs it in
    a worker thread instead of stalling the event loop for every other request.

    No auth gate -- see the module docstring's AUTH NOTE.
    """
    if _graph is None:
        raise HTTPException(status_code=503, detail="Agent is still starting up.")

    started = time.perf_counter()
    log.info("research query: %r", req.query)
    try:
        final = _graph.invoke(_initial_state(req.query), config={"recursion_limit": 25})
    except Exception as exc:  # noqa: BLE001
        log.exception("agent invocation failed")
        raise HTTPException(status_code=502, detail=f"Agent failed: {exc}") from exc

    result = final.get("result")
    if result is None:
        raise HTTPException(status_code=502, detail="Agent produced no structured result.")

    log.info("done in %.1fs — %d rows", time.perf_counter() - started, len(result.table_data))
    return result


# --- SSE node -> friendly status label ---------------------------------------
# extract_trial is deliberately absent here -- it is reported via `progress`
# events (one status line per Map worker would flood the stream when a
# broad query fans out to dozens or hundreds of them), not `status`.
_NODE_STATUS_LABELS = {
    "intent_classifier": "IntentClassifier",
    "agent": "Agent",
    "tools": "ToolNode",
    "synthesize_table": "Reducer",
    "out_of_domain": "OutOfDomain",
    "no_results_fallback": "NoResultsFallback",
}


def _sse(event: str, data: dict[str, Any]) -> str:
    """One well-formed SSE chunk. `data` must already be JSON-serialisable
    (see _jsonable below for LangChain/Pydantic objects that aren't, by
    default)."""
    return f"event: {event}\ndata: {json.dumps(data)}\n\n"


def _jsonable(value: Any) -> Any:
    """LangGraph task_result payloads carry LangChain message objects and,
    at the terminal node, a SmartTableResponse -- neither is JSON-serialisable
    as-is. model_dump()/model_dump_json() handles the Pydantic case; anything
    else (e.g. a list of AIMessage/ToolMessage) falls back to str() so a
    single odd value can never crash the whole stream over one field the
    frontend was never going to render anyway."""
    if isinstance(value, BaseModel):
        return json.loads(value.model_dump_json())
    if isinstance(value, list):
        return [_jsonable(v) for v in value]
    if isinstance(value, dict):
        return {k: _jsonable(v) for k, v in value.items()}
    if isinstance(value, (str, int, float, bool)) or value is None:
        return value
    return str(value)


async def _stream_events(query: str, tenant_id: str | None = None) -> AsyncIterator[str]:
    """The actual SSE generator -- translates LangGraph's stream_mode="debug"
    events (verified directly against this exact graph before writing this:
    a "task" event fires BEFORE a node runs, a "task_result" event fires
    AFTER, each payload carrying the node's `name` and, for task_result, its
    `result` state-update dict) into the three event types the spec calls
    for.

    Termination is detected generically, not by hardcoding which node name
    is "the last one": AgentState.result is set ONLY by a node representing
    a genuine terminal outcome (synthesize_table on a validated response,
    out_of_domain, no_results_fallback) -- never by an in-progress synthesis
    retry. Checking for a non-None `result` key in ANY task_result's own
    result dict is therefore a correct, node-name-agnostic way to know the
    graph is done, and it naturally covers all three exit paths with one
    check instead of three.
    """
    started = time.perf_counter()
    map_total = 0
    map_done = 0
    map_started_announced = False

    try:
        async for event in _graph.astream(
            _initial_state(query, tenant_id=tenant_id), config={"recursion_limit": 25},
            stream_mode="debug"
        ):
            etype = event.get("type")
            payload = event.get("payload") or {}
            name = payload.get("name")

            if etype == "task":
                if name == "extract_trial":
                    # Dispatched together as one Send-based fan-out per
                    # LangGraph superstep -- all "task" events for this
                    # step arrive before any of their "task_result"s, so
                    # counting them here gives an accurate denominator for
                    # the progress events below.
                    map_total += 1
                    continue
                label = _NODE_STATUS_LABELS.get(name, name)
                yield _sse("status", {"node": label, "phase": "start"})
                continue

            if etype != "task_result":
                continue

            if name == "extract_trial":
                map_done += 1
                if not map_started_announced:
                    yield _sse("status", {"node": "MapWorkers", "phase": "start",
                                          "total": map_total})
                    map_started_announced = True
                yield _sse("progress", {"node": "MapWorkers", "completed": map_done,
                                        "total": map_total})
                continue

            label = _NODE_STATUS_LABELS.get(name, name)
            if payload.get("error"):
                yield _sse("status", {"node": label, "phase": "error",
                                      "error": str(payload["error"])})
                continue

            yield _sse("status", {"node": label, "phase": "done"})

            result_update = payload.get("result") or {}
            final_result = result_update.get("result")
            if final_result is not None:
                yield _sse("result", _jsonable(final_result))
                log.info("stream done in %.1fs — %d rows", time.perf_counter() - started,
                         len(final_result.table_data))
                return

        # The graph ran to completion without any node ever setting
        # AgentState.result -- a genuine bug in the graph, not a network
        # hiccup, but the stream must still end with SOMETHING the frontend
        # can act on rather than just closing silently.
        yield _sse("error", {"message": "Agent stream ended without producing a result."})

    except Exception as exc:  # noqa: BLE001 -- must reach the client as an SSE event, not a raw 500
        log.exception("stream failed")
        yield _sse("error", {"message": str(exc)})


@app.post("/api/research/stream")
async def research_stream(req: ResearchRequest) -> StreamingResponse:
    """Same agent, same request shape as POST /api/research, but responds
    with a live text/event-stream instead of waiting for the whole run to
    finish. See module docstring for why this exists.

    A plain browser EventSource can't be used to CONSUME this from the
    frontend -- EventSource only supports GET, and this is a POST (to keep
    the same ResearchRequest body /api/research already uses, and to avoid
    putting an analyst's free-text question in a URL). The frontend instead
    reads response.body via fetch()'s own ReadableStream, which works for
    any method.

    No auth gate -- see the module docstring's AUTH NOTE.
    """
    if _graph is None:
        raise HTTPException(status_code=503, detail="Agent is still starting up.")

    log.info("research stream query: %r", req.query)
    return StreamingResponse(
        _stream_events(req.query),
        media_type="text/event-stream",
        headers={
            # Disables buffering on nginx-style reverse proxies that would
            # otherwise hold the whole response until it closes -- exactly
            # the failure mode this endpoint exists to avoid at the ALB
            # layer too.
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# =============================================================================
# INDICATION LANDSCAPE -- a single-shot competitive matrix (mechanism/target
# rows x development-phase columns) for a therapeutic area, via the separate
# retrieve-then-synthesize graph in research_agent.py. See that module's
# LandscapeMatrix/make_landscape_graph docstrings for why this is a distinct
# pipeline from POST /api/research rather than a variant of it.
# =============================================================================
class LandscapeRequest(BaseModel):
    therapeutic_area: str = Field(
        min_length=3,
        max_length=200,
        description="A therapeutic area / indication to build a competitive "
                    "landscape matrix for.",
        examples=["Non-Small Cell Lung Cancer"],
    )


@app.post("/api/landscape", response_model=LandscapeMatrix)
def landscape(req: LandscapeRequest) -> LandscapeMatrix:
    """Run the landscape graph and return the validated competitive matrix.

    Declared `def`, not `async def`, for the same reason as POST
    /api/research: graph.invoke() is blocking (Qdrant retrieval across
    three collections plus one LLM call), so FastAPI runs it in a worker
    thread instead of stalling the event loop.

    No auth gate -- see the module docstring's AUTH NOTE.
    """
    if _landscape_graph is None:
        raise HTTPException(status_code=503, detail="Landscape agent is still starting up.")

    started = time.perf_counter()
    log.info("landscape query: %r", req.therapeutic_area)
    try:
        result = run_landscape_query(_landscape_graph, req.therapeutic_area)
    except Exception as exc:  # noqa: BLE001
        log.exception("landscape agent invocation failed")
        raise HTTPException(status_code=502, detail=f"Landscape agent failed: {exc}") from exc

    if result is None:
        raise HTTPException(status_code=502, detail="Landscape agent produced no structured result.")

    log.info("landscape done in %.1fs — %d mechanism row(s)",
             time.perf_counter() - started, len(result.rows))
    return result


# =============================================================================
# CATALYST & READOUT TRACKER -- a chronological timeline of upcoming market-
# moving events for a query, via yet another retrieve-then-synthesize graph.
# See research_agent.py's CatalystTimeline/make_catalyst_graph docstrings
# for the live ClinicalTrials.gov date-enrichment step this one does that
# neither /api/research nor /api/landscape need.
# =============================================================================
class CatalystRequest(BaseModel):
    query: str = Field(
        min_length=3,
        max_length=200,
        description="A therapeutic area and/or event type to build a "
                    "catalyst timeline for.",
        examples=["Upcoming Phase 3 readouts in Oncology"],
    )


@app.post("/api/catalysts", response_model=CatalystTimeline)
def catalysts(req: CatalystRequest) -> CatalystTimeline:
    """Run the catalyst graph and return the validated chronological timeline.

    Declared `def`, not `async def`, for the same reason as POST
    /api/research and /api/landscape: graph.invoke() is blocking (Qdrant
    retrieval, a live ClinicalTrials.gov date lookup, and one LLM call), so
    FastAPI runs it in a worker thread instead of stalling the event loop.

    No auth gate -- see the module docstring's AUTH NOTE.
    """
    if _catalyst_graph is None:
        raise HTTPException(status_code=503, detail="Catalyst agent is still starting up.")

    started = time.perf_counter()
    log.info("catalyst query: %r", req.query)
    try:
        result = run_catalyst_query(_catalyst_graph, req.query)
    except Exception as exc:  # noqa: BLE001
        log.exception("catalyst agent invocation failed")
        raise HTTPException(status_code=502, detail=f"Catalyst agent failed: {exc}") from exc

    if result is None:
        raise HTTPException(status_code=502, detail="Catalyst agent produced no structured result.")

    log.info("catalyst done in %.1fs — %d event(s)",
             time.perf_counter() - started, len(result.events))
    return result


# =============================================================================
# EXECUTIVE EXPORT -- POST a SmartTableResponse the frontend already has in
# state (from either /api/research or /api/research/stream's `result` event)
# back to the server to render as a downloadable file. Deliberately NOT
# re-running the agent: the export is of whatever the analyst is already
# looking at, so the request body IS the SmartTableResponse itself, not a
# query string -- these endpoints do no LLM/Qdrant/Neo4j work at all.
# =============================================================================
EXCEL_HEADERS = ["NCT ID", "Sponsor", "Phase", "Interventions", "Mechanism / Findings",
                 "Mechanism Described"]
# Column widths tuned for this schema's actual content shape (mechanism text
# runs long, phase/sponsor are short) -- not left at openpyxl's default,
# which would make the export technically correct but unreadable without
# the analyst manually resizing every column first.
EXCEL_COLUMN_WIDTHS = [14, 28, 16, 34, 70, 18]
EXCEL_HEADER_FILL = "1E3A5F"  # matches the frontend's sky-900-ish header tone


@app.post("/api/export/excel")
def export_excel(data: SmartTableResponse) -> Response:
    """table_data -> one formatted sheet, .xlsx. Fully in-memory (io.BytesIO)
    -- these exports are at most a few hundred rows, nowhere near large
    enough to justify streaming the workbook to disk first.

    No auth gate -- see the module docstring's AUTH NOTE.
    """
    log.info("export xlsx: %d rows", len(data.table_data))
    wb = Workbook()
    ws = wb.active
    ws.title = "Clinical Trials"

    ws.append(EXCEL_HEADERS)
    header_fill = PatternFill(start_color=EXCEL_HEADER_FILL, end_color=EXCEL_HEADER_FILL,
                              fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font

    for row in data.table_data:
        ws.append([
            row.nct_id,
            row.sponsor,
            row.phase,
            ", ".join(row.interventions),
            row.mechanism_or_findings,
            "Yes" if row.mechanism_described else "No",
        ])

    for i, width in enumerate(EXCEL_COLUMN_WIDTHS, start=1):
        ws.column_dimensions[get_column_letter(i)].width = width
    ws.freeze_panes = "A2"  # header stays visible when scrolling a long export

    buf = io.BytesIO()
    wb.save(buf)

    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": 'attachment; filename="clinical_trials_export.xlsx"'},
    )


PPTX_COLUMNS = ["NCT ID", "Sponsor", "Phase", "Interventions", "Mechanism / Findings"]
# Verified directly against a live render (not guessed): 8 data rows + 1
# header row on a Title-Only layout's default content area (0.3"-9.7"
# wide, 1.3"-6.8" tall on a standard 10x7.5" slide) keeps each cell's text
# legible at Pt(9) without a table so dense it becomes a wall of
# unreadable text -- the actual failure mode a single giant table across
# hundreds of rows would hit.
PPTX_ROWS_PER_SLIDE = 8


def _set_pptx_header_cell(cell, text: str) -> None:
    cell.text = text
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.size = Pt(11)


def _set_pptx_data_cell(cell, text: str) -> None:
    cell.text = text
    cell.text_frame.paragraphs[0].font.size = Pt(9)


@app.post("/api/export/pptx")
def export_pptx(data: SmartTableResponse) -> Response:
    """Title slide -> briefing slide -> one or more data-table slides
    (chunked at PPTX_ROWS_PER_SLIDE so a large result set doesn't collapse
    into one unreadable mega-table), .pptx.

    No auth gate -- see the module docstring's AUTH NOTE.
    """
    log.info("export pptx: %d rows", len(data.table_data))
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # "Title Slide"
    title_slide.shapes.title.text = "Clinical Landscape Analysis"
    if len(title_slide.placeholders) > 1:
        n = len(data.table_data)
        title_slide.placeholders[1].text = f"{n} trial{'s' if n != 1 else ''} analyzed"

    briefing_slide = prs.slides.add_slide(prs.slide_layouts[1])  # "Title and Content"
    briefing_slide.shapes.title.text = "Briefing"
    body = briefing_slide.placeholders[1].text_frame
    body.text = data.narrative_summary
    body.word_wrap = True

    rows = data.table_data
    # range(0, 0, N) is simply empty -- an explicit "no trials" slide reads
    # better to an analyst than a presentation that silently has no data
    # section at all.
    chunk_starts = range(0, len(rows), PPTX_ROWS_PER_SLIDE) if rows else [0]
    for chunk_start in chunk_starts:
        chunk = rows[chunk_start:chunk_start + PPTX_ROWS_PER_SLIDE]
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only"
        slide.shapes.title.text = (
            f"Trial Data ({chunk_start + 1}–{chunk_start + len(chunk)} of {len(rows)})"
            if chunk else "Trial Data (no trials in this result)"
        )

        n_rows = len(chunk) + 1  # +1 header row; still valid (2 rows) when chunk is empty
        table = slide.shapes.add_table(
            max(n_rows, 2), len(PPTX_COLUMNS), Inches(0.3), Inches(1.3), Inches(9.4), Inches(5.5)
        ).table

        for c, col_name in enumerate(PPTX_COLUMNS):
            _set_pptx_header_cell(table.cell(0, c), col_name)

        for r, row in enumerate(chunk, start=1):
            values = [row.nct_id, row.sponsor, row.phase,
                     ", ".join(row.interventions), row.mechanism_or_findings]
            for c, val in enumerate(values):
                _set_pptx_data_cell(table.cell(r, c), val)

    buf = io.BytesIO()
    prs.save(buf)

    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="clinical_landscape_analysis.pptx"'},
    )
